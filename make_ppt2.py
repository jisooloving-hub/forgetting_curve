from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 폰트 ───────────────────────────────────
FT  = "NanumGothicExtraBold"   # 제목
FB  = "Noto Sans KR Light"     # 본문

# ── 색상 ───────────────────────────────────
C_BG    = RGBColor(0x0F, 0x0F, 0x1A)   # 거의 검정
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE  = RGBColor(0x4A, 0x90, 0xD9)
C_PURP  = RGBColor(0x9B, 0x59, 0xB6)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_ORAN  = RGBColor(0xE6, 0x7E, 0x22)
C_RED   = RGBColor(0xE7, 0x4C, 0x3C)
C_DARK  = RGBColor(0x1E, 0x1E, 0x2E)
C_GRAY  = RGBColor(0xF4, 0xF4, 0xF8)
C_LGRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_DONE  = RGBColor(0x27, 0xAE, 0x60)
C_WIP   = RGBColor(0xF3, 0x9C, 0x12)
C_TODO  = RGBColor(0xCC, 0xCC, 0xCC)

def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def box(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=16, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, font=None, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name  = font or (FT if bold else FB)
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    return tb

def img(slide, path, x, y, w, h):
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

# ════════════════════════════════════════════
# 슬라이드 1 — 타이틀
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)

# 왼쪽 파란 세로 바
box(sl, 0, 0, 0.12, 7.5, C_BLUE)

# 앱 이름
txt(sl, "커브터디", 0.5, 1.8, 9, 1.4, size=60, bold=True, color=C_WHITE, font=FT)
txt(sl, "curvetudy", 0.6, 3.1, 6, 0.6, size=22, color=C_BLUE, font=FB)

# 구분선
box(sl, 0.5, 3.85, 5.5, 0.05, C_BLUE)

# 부제목
txt(sl, "망각곡선 이론 기반  ·  복습 타이밍 자동 추천 앱",
    0.5, 4.05, 10, 0.6, size=17, color=C_LGRAY, font=FB)
txt(sl, "React Native (Expo)  ·  Firebase  ·  Claude API",
    0.5, 4.65, 10, 0.5, size=14, color=RGBColor(0x66,0x66,0x88), font=FB)

txt(sl, "2026. 05", 0.5, 6.6, 5, 0.4, size=13,
    color=RGBColor(0x55,0x55,0x66), font=FB)

# ════════════════════════════════════════════
# 슬라이드 2 — 문제 정의 & 비전
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_BG)

# 상단 레이블
txt(sl, "PROBLEM  →  SOLUTION", 0.5, 0.28, 12, 0.4,
    size=11, color=C_BLUE, font=FB)
box(sl, 0.5, 0.65, 12.3, 0.04, C_BLUE)

# ── 왼쪽: 문제 ──────────────────────
txt(sl, "문제 정의", 0.5, 0.85, 5.5, 0.5,
    size=13, color=C_RED, bold=True, font=FT)

txt(sl, "왜 우리는\n공부해도 잊을까?",
    0.5, 1.35, 5.8, 1.5, size=34, bold=True, color=C_WHITE, font=FT)

txt(sl, "에빙하우스(Ebbinghaus, 1885)의 연구에 따르면\n"
        "사람은 학습 후 시간이 지날수록 기억이 급격히 감소합니다.\n"
        "반복 복습 없이는 장기 기억으로 이어지지 않습니다.",
    0.5, 2.9, 5.8, 1.2, size=13, color=C_LGRAY, font=FB)

# 망각률 수치 박스 3개
stats = [("20분 후", "42%", "망각"), ("1일 후", "70%", "망각"), ("1주일 후", "77%", "망각")]
for i, (label, pct, sub) in enumerate(stats):
    bx = 0.5 + i * 2.0
    box(sl, bx, 4.25, 1.7, 1.8, RGBColor(0x2A, 0x10, 0x10))
    txt(sl, label, bx, 4.35, 1.7, 0.35, size=10, color=C_LGRAY,
        align=PP_ALIGN.CENTER, font=FB)
    txt(sl, pct,   bx, 4.65, 1.7, 0.65, size=30, bold=True, color=C_RED,
        align=PP_ALIGN.CENTER, font=FT)
    txt(sl, sub,   bx, 5.25, 1.7, 0.35, size=11, color=C_LGRAY,
        align=PP_ALIGN.CENTER, font=FB)

# 중앙 화살표
txt(sl, "→", 6.55, 3.2, 0.8, 0.8, size=36, color=C_BLUE,
    align=PP_ALIGN.CENTER, font=FT)

# ── 오른쪽: 해결 ──────────────────────
box(sl, 6.9, 0.85, 6.0, 5.5, RGBColor(0x0A, 0x18, 0x2E))

txt(sl, "해결책", 7.2, 1.0, 5, 0.45,
    size=13, color=C_BLUE, bold=True, font=FT)
txt(sl, "커브터디", 7.2, 1.45, 5, 0.9,
    size=36, bold=True, color=C_WHITE, font=FT)

solutions = [
    ("📝", "노트 작성", "공부한 내용을 노트로 기록"),
    ("🤖", "자동 계산", "망각곡선 기반 최적 복습 시점 계산"),
    ("✅", "자가 체크", "잘됨 → 3일 후  /  안됨 → 1일 후"),
    ("🤖", "AI 문제 생성", "Claude API로 노트 기반 퀴즈 자동 생성"),
]
for i, (icon, title, desc) in enumerate(solutions):
    y = 2.5 + i * 0.9
    txt(sl, icon,  7.2,  y,       0.5, 0.6, size=16, font=FB)
    txt(sl, title, 7.7,  y,       2.2, 0.35, size=14, bold=True,
        color=C_WHITE, font=FT)
    txt(sl, desc,  7.7,  y+0.35,  5.0, 0.4,  size=11,
        color=C_LGRAY, font=FB)

# ════════════════════════════════════════════
# 슬라이드 3 — WBS
# ════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_GRAY)

# 제목 바
box(sl, 0, 0, 13.33, 1.0, C_DARK)
txt(sl, "WBS — 프로젝트 작업 분류", 0.4, 0.1, 10, 0.5,
    size=22, bold=True, color=C_WHITE, font=FT)
txt(sl, "Work Breakdown Structure  |  커브터디 전체 개발 현황",
    0.4, 0.58, 10, 0.35, size=12, color=C_LGRAY, font=FB)

# 범례
for lx, lc, ll in [(10.5, C_DONE, "완료"), (11.3, C_WIP, "진행 중"), (12.2, C_TODO, "예정")]:
    box(sl, lx, 0.28, 0.18, 0.18, lc)
    txt(sl, ll, lx+0.22, 0.25, 0.8, 0.3, size=10, color=C_WHITE, font=FB)

# WBS 카테고리 정의
categories = [
    ("기획",    C_BLUE,  [("1.1", "앱 기획 (ARD 작성)",         "완료"),
                           ("1.2", "기술 스택 선택",               "완료")]),
    ("설계",    C_PURP,  [("2.1", "시스템 아키텍처",              "완료"),
                           ("2.2", "DB 스키마 설계",               "완료"),
                           ("2.3", "화면(UI) 설계",                "완료")]),
    ("개발",    C_GREEN, [("3.1", "프로젝트 세팅",                 "완료"),
                           ("3.2", "인증 (로그인/회원가입)",        "완료"),
                           ("3.3", "핵심 기능 (노트·복습)",         "진행중"),
                           ("3.4", "캘린더 + 스트릭",              "예정"),
                           ("3.5", "Claude API 연동",              "예정"),
                           ("3.6", "UI 마무리·설정",               "예정")]),
    ("테스트",  C_ORAN,  [("4.1", "기능 테스트",                   "예정"),
                           ("4.2", "UI/UX 테스트",                 "예정")]),
    ("배포",    C_RED,   [("5.1", "Vercel 배포",                   "예정")]),
]

STATUS_COLOR = {"완료": C_DONE, "진행중": C_WIP, "예정": C_TODO}
STATUS_LABEL = {"완료": "✅ 완료", "진행중": "🔄 진행 중", "예정": "⬜ 예정"}

col_w   = 2.3
col_gap = 0.18
start_x = (13.33 - (col_w * 5 + col_gap * 4)) / 2
top_y   = 1.15

for ci, (cat_name, cat_color, items) in enumerate(categories):
    cx = start_x + ci * (col_w + col_gap)

    # 카테고리 헤더
    hdr = box(sl, cx, top_y, col_w, 0.55, cat_color)
    txt(sl, cat_name, cx, top_y + 0.08, col_w, 0.4,
        size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font=FT)

    # 항목들
    for ri, (code, name, status) in enumerate(items):
        iy = top_y + 0.55 + 0.02 + ri * 0.82
        # 카드 배경
        box(sl, cx, iy, col_w, 0.76, C_WHITE, line=RGBColor(0xDD,0xDD,0xEE))
        # 색상 왼쪽 바
        box(sl, cx, iy, 0.07, 0.76, cat_color)
        # 코드 + 이름
        txt(sl, code, cx+0.12, iy+0.04, col_w-0.15, 0.24,
            size=9, color=C_LGRAY, font=FB)
        txt(sl, name, cx+0.12, iy+0.24, col_w-0.15, 0.28,
            size=11, bold=False, color=C_DARK, font=FB)
        # 상태 뱃지
        sc = STATUS_COLOR[status]
        box(sl, cx+0.12, iy+0.5, col_w-0.25, 0.2, sc)
        txt(sl, STATUS_LABEL[status], cx+0.12, iy+0.5, col_w-0.25, 0.22,
            size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font=FB)

# ── 다이어그램 슬라이드 공통 생성 함수 ──────
DIAG_DIR = r"C:\Users\jisoo\OneDrive\Desktop\클로드 바이브 코딩\바이브 코딩 개인 노트\망각곡선\diagrams"

def make_diagram_slide(title, subtitle, filename, img_w=12.1, img_h=6.1):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, C_WHITE)
    box(sl, 0, 0, 13.33, 1.0, C_DARK)
    txt(sl, title,    0.4, 0.1,  10, 0.5,  size=22, bold=True, color=C_WHITE, font=FT)
    txt(sl, subtitle, 0.4, 0.58, 10, 0.35, size=12, color=C_LGRAY, font=FB)
    # 이미지 가로 중앙 정렬
    img_x = (13.33 - img_w) / 2
    img_y = (7.5 - img_h) / 2 + 0.5
    try:
        img(sl, rf"{DIAG_DIR}\{filename}", img_x, img_y, img_w, img_h)
    except Exception as e:
        txt(sl, f"[이미지 없음: {e}]", 1, 3, 11, 1, size=14, color=C_RED, font=FB)

# ════════════════════════════════════════════
# 슬라이드 4 — 시스템 아키텍처
# ════════════════════════════════════════════
make_diagram_slide(
    "시스템 아키텍처",
    "Expo React Native  ·  Firebase  ·  Claude API 연동 구조",
    "01_architecture.png"
)

# ════════════════════════════════════════════
# 슬라이드 5 — 화면 네비게이션
# ════════════════════════════════════════════
make_diagram_slide(
    "화면 네비게이션 흐름",
    "로그인 → 홈 → 노트 작성 / 복습 / 전체 목록 전체 화면 이동 경로",
    "02_navigation.png"
)

# ════════════════════════════════════════════
# 슬라이드 6 — 데이터베이스 구조
# ════════════════════════════════════════════
make_diagram_slide(
    "데이터베이스 구조",
    "Firestore  ·  users / notes / reviews 컬렉션 관계도",
    "03_database.png",
    img_w=5.0, img_h=6.1
)

# ── 저장 ──────────────────────────────────
output = r"C:\Users\jisoo\OneDrive\Desktop\클로드 바이브 코딩\바이브 코딩 개인 노트\망각곡선\커브터디_발표_v2.pptx"
prs.save(output)
print("저장 완료: " + output)
