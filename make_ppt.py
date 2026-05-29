from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 폰트 정의 ──────────────────────────────
FONT_TITLE = "NanumGothicExtraBold"  # 제목용 — 굵고 임팩트
FONT_BODY  = "Noto Sans KR Light"    # 본문용 — 얇고 세련됨

# ── 색상 정의 ──────────────────────────────
BLACK   = RGBColor(0x1A, 0x1A, 0x2E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x4A, 0x90, 0xD9)
PURPLE  = RGBColor(0x7B, 0x68, 0xEE)
GREEN   = RGBColor(0x2E, 0xCC, 0x71)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
GRAY    = RGBColor(0xF5, 0xF5, 0xF5)
DARK    = RGBColor(0x2C, 0x3E, 0x50)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_TITLE if bold else FONT_BODY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_rect(slide, x, y, w, h, fill_color, radius=False):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_card(slide, x, y, w, h, title, title_color, content_lines):
    """카드 박스: 색상 헤더 + 흰 본문"""
    # 헤더
    header = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.55))
    header.fill.solid()
    header.fill.fore_color.rgb = title_color
    header.line.fill.background()
    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # 본문
    body = slide.shapes.add_shape(1, Inches(x), Inches(y+0.55), Inches(w), Inches(h-0.55))
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    tf2 = body.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(content_lines):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = line
        run2.font.name = FONT_BODY
        run2.font.size = Pt(12)
        run2.font.color.rgb = DARK

# ════════════════════════════════════════
# 슬라이드 1 — 타이틀
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BLACK)

add_rect(slide, 0, 2.8, 13.33, 0.08, BLUE)

add_textbox(slide, "🧠 망각곡선 앱", 1.5, 1.2, 10, 1.2,
            font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "복습 타이밍 자동 추천 앱 — React Native 프로토타입 개발",
            1.5, 2.5, 10, 0.6,
            font_size=20, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(slide, "2026. 05. 25",
            1.5, 5.8, 10, 0.6,
            font_size=14, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# 슬라이드 2 — 기술 스택
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)

# 상단 제목 바
bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
bar.fill.solid()
bar.fill.fore_color.rgb = DARK
bar.line.fill.background()
tf = bar.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "  기술 스택 선택"
run.font.name = FONT_TITLE
run.font.size = Pt(26)
run.font.bold = True
run.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "  각 기술을 선택한 이유"
r2.font.name = FONT_BODY
r2.font.size = Pt(13)
r2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)

stacks = [
    (BLUE,   "Expo\n(React Native)",
     "앱 프레임워크",
     [
         "선택 이유:",
         "• Windows에서도 개발 가능",
         "  (React Native CLI는",
         "   Mac 없으면 iOS 불가)",
         "• Expo Go 앱으로 QR 코드",
         "  스캔 → 즉시 폰 테스트",
         "• 설정 5분 만에 시작 가능",
     ]),
    (RGBColor(0xFF, 0x9A, 0x00), "Firebase",
     "인증 + 데이터베이스",
     [
         "선택 이유:",
         "• 직접 서버 안 만들어도 됨",
         "• Auth: 이메일 로그인",
         "  코드 5줄로 구현",
         "• Firestore: NoSQL DB",
         "  실시간 데이터 동기화",
         "• 무료 플랜(Spark)으로",
         "  학교 프로젝트 수준 충분",
     ]),
    (PURPLE, "Expo Router",
     "화면 라우팅",
     [
         "선택 이유:",
         "• 파일 이름 = 화면 경로",
         "  login.tsx → /login",
         "• 폴더 구조만 봐도",
         "  앱 구조 파악 가능",
         "• Next.js와 유사한 방식",
         "  → 웹 개발 경험 활용",
     ]),
    (GREEN,  "Claude API",
     "AI 문제 생성",
     [
         "선택 이유:",
         "• 노트 내용을 넘기면",
         "  문제 자동 생성",
         "• 앱 차별화 포인트",
         "  (단순 자가 체크 대비)",
         "• 무료 크레딧으로 시작",
         "• 발표 시 'AI 활용'",
         "  요소로 어필 가능",
     ]),
    (RGBColor(0x23, 0x9B, 0xD3), "TypeScript",
     "개발 언어",
     [
         "선택 이유:",
         "• 타입 정의로 데이터",
         "  구조 명확히 표현",
         "• 오타/타입 오류를",
         "  실행 전에 잡아줌",
         "• 자동완성 지원으로",
         "  개발 속도 향상",
     ]),
]

card_w = 2.3
card_h = 5.3
gap = 0.28
start_x = (13.33 - (card_w * 5 + gap * 4)) / 2

for i, (color, name, role, lines) in enumerate(stacks):
    x = start_x + i * (card_w + gap)
    y = 1.2

    # 색상 헤더 (이름 + 역할)
    header = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(card_w), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = name
    run.font.name = FONT_TITLE
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = role
    r2.font.name = FONT_BODY
    r2.font.size = Pt(10)
    r2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xFF)

    # 본문
    body = slide.shapes.add_shape(1, Inches(x), Inches(y + 0.9), Inches(card_w), Inches(card_h - 0.9))
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    body.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    tf2 = body.text_frame
    tf2.word_wrap = True
    for j, line in enumerate(lines):
        p3 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        run3 = p3.add_run()
        run3.text = line
        run3.font.name = FONT_BODY
        run3.font.size = Pt(11)
        run3.font.bold = (line == "선택 이유:")
        run3.font.color.rgb = DARK

# ════════════════════════════════════════
# 슬라이드 4 — 전체 개발 단계
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, GRAY)

add_textbox(slide, "전체 개발 계획 — 6단계", 0.5, 0.3, 12, 0.7,
            font_size=28, bold=True, color=DARK)
add_textbox(slide, "프로토타입 목표: 1~4단계 완성 후 시연 가능한 수준",
            0.5, 0.95, 12, 0.4, font_size=14, color=RGBColor(0x66,0x66,0x66))

stages = [
    (BLUE,   "1단계\n프로젝트 세팅",   "✅ 완료"),
    (PURPLE, "2단계\n인증",             "✅ 완료"),
    (GREEN,  "3단계\n노트 + 복습",      "🔄 진행 중"),
    (ORANGE, "4단계\n캘린더 + 스트릭", "⬜ 예정"),
    (RGBColor(0xE7,0x4C,0x3C), "5단계\nAI 테스트", "⬜ 예정"),
    (RGBColor(0x95,0xA5,0xA6), "6단계\nUI 마무리", "⬜ 예정"),
]

for i, (color, label, status) in enumerate(stages):
    x = 0.5 + i * 2.05
    box = slide.shapes.add_shape(1, Inches(x), Inches(1.6), Inches(1.85), Inches(3.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT_TITLE
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "\n" + status
    r2.font.name = FONT_BODY
    r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE

add_textbox(slide, "※ 현재 3단계 진행 중", 0.5, 5.1, 12, 0.4,
            font_size=13, color=RGBColor(0x88,0x88,0x88))

# ── 공통 슬라이드 생성 함수 ───────────────
def make_stage_slide(title, subtitle, cards):
    """
    cards = [
      ("핵심 개념", BLUE,   ["줄1","줄2",...]),
      ("프로젝트 이유", PURPLE, [...]),
      ("오늘 산출물", GREEN,  [...]),
      ("저장 위치", ORANGE,  [...]),
    ]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    # 상단 제목 바
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tf = bar.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + title
    run.font.name = FONT_TITLE
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "  " + subtitle
    r2.font.name = FONT_BODY
    r2.font.size = Pt(13)
    r2.font.color.rgb = RGBColor(0xAA,0xCC,0xFF)

    # 카드 4개
    card_w = 2.9
    card_h = 5.2
    gap = 0.38
    start_x = (13.33 - (card_w * 4 + gap * 3)) / 2
    for i, (card_title, card_color, lines) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, 1.3, card_w, card_h, card_title, card_color, lines)

    return slide

# ════════════════════════════════════════
# 슬라이드 3 — 앱 기획 (ARD)
# ════════════════════════════════════════
make_stage_slide(
    "앱 기획 — ARD 작성",
    "무엇을, 왜, 어떻게 만들 것인가를 문서로 정리",
    [
        ("핵심 개념", BLUE, [
            "망각곡선 이론 기반",
            "복습 타이밍 자동 추천 앱",
            "",
            "• 공부한 내용을 노트로 저장",
            "• 이해도 테스트 후 복습 날짜",
            "  자동 계산",
            "• 잘됨 → 3일 후",
            "• 안됨 → 1일 후",
        ]),
        ("프로젝트 이유", PURPLE, [
            "언제 복습해야 할지",
            "스스로 판단하기 어렵다",
            "",
            "• 너무 빨리 → 시간 낭비",
            "• 너무 늦게 → 이미 망각",
            "",
            "최적 복습 시점을",
            "앱이 자동으로 계산",
        ]),
        ("오늘 산출물", GREEN, [
            "ARD.md 작성 완료",
            "",
            "• 앱 목적 & 타겟",
            "• 주요 기능 5가지",
            "• 사용자 흐름 4가지",
            "• 화면 목록 8개",
            "• 데이터 구조 정의",
        ]),
        ("저장 위치", ORANGE, [
            "📁 망각곡선/",
            "└── ARD.md",
            "",
            "앱 요구사항 문서",
            "(App Requirements Doc)",
        ]),
    ]
)

# ════════════════════════════════════════
# 슬라이드 4 — 1단계: 프로젝트 세팅
# ════════════════════════════════════════
make_stage_slide(
    "1단계 — 프로젝트 세팅",
    "개발 환경 구성 및 폴더 구조 설계",
    [
        ("핵심 개념", BLUE, [
            "Expo + TypeScript",
            "프로젝트 초기화",
            "",
            "• Expo Router 기반",
            "  파일 = 화면 구조",
            "• 폴더별 역할 분리",
            "  app / services /",
            "  utils / types",
        ]),
        ("프로젝트 이유", PURPLE, [
            "빠른 시작 & 유지보수",
            "용이한 구조 필요",
            "",
            "• Expo: 윈도우에서도",
            "  iOS/Android 개발 가능",
            "• 폴더 분리: 화면과",
            "  데이터 로직 분리로",
            "  코드 가독성 향상",
        ]),
        ("오늘 산출물", GREEN, [
            "프로젝트 뼈대 완성",
            "",
            "• forgetting-curve/",
            "  프로젝트 생성",
            "• 6개 라이브러리 설치",
            "• Expo Router 세팅",
            "• 폴더 구조 생성",
        ]),
        ("저장 위치", ORANGE, [
            "📁 forgetting-curve/",
            "├── package.json",
            "├── app.json",
            "└── app/_layout.tsx",
            "",
            "개발일지.md",
            "1단계 섹션",
        ]),
    ]
)

# ════════════════════════════════════════
# 슬라이드 5 — 2단계: 인증
# ════════════════════════════════════════
make_stage_slide(
    "2단계 — 인증 구현",
    "Firebase Auth 기반 이메일 로그인 / 회원가입",
    [
        ("핵심 개념", BLUE, [
            "Firebase Authentication",
            "이메일 로그인",
            "",
            "• 로그인 / 회원가입 화면",
            "• 로그인 상태 유지",
            "  (앱 재시작 후에도)",
            "• 미인증 시 자동으로",
            "  로그인 화면 이동",
        ]),
        ("프로젝트 이유", PURPLE, [
            "사용자별 데이터 분리",
            "필요",
            "",
            "• 노트 / 복습 기록이",
            "  본인 것만 보여야 함",
            "• 직접 서버 없이",
            "  Firebase로 5줄 구현",
        ]),
        ("오늘 산출물", GREEN, [
            "로그인 화면 완성",
            "",
            "• login.tsx",
            "• signup.tsx",
            "• services/auth.ts",
            "• 유효성 검사 포함",
            "  (빈칸, 비번 6자 이상,",
            "   비번 확인 일치)",
        ]),
        ("저장 위치", ORANGE, [
            "📁 app/(auth)/",
            "├── _layout.tsx",
            "├── login.tsx",
            "└── signup.tsx",
            "",
            "📁 services/",
            "└── auth.ts",
        ]),
    ]
)

# ════════════════════════════════════════
# 슬라이드 6 — 3단계: 핵심 기능
# ════════════════════════════════════════
make_stage_slide(
    "3단계 — 핵심 기능 (진행 중)",
    "노트 작성 + 홈 화면 (오늘 복습 목록)",
    [
        ("핵심 개념", BLUE, [
            "노트 저장 &",
            "복습 일정 자동 계산",
            "",
            "• 노트 작성 시",
            "  1일 후 복습 자동 설정",
            "• 홈 화면에서",
            "  오늘 복습 목록 표시",
            "• Firestore에 저장",
        ]),
        ("프로젝트 이유", PURPLE, [
            "앱의 핵심 가치 구현",
            "",
            "• '언제 복습해야 하나'",
            "  고민 없이 자동 계산",
            "• 노트 없이는 복습도",
            "  없으므로 가장 먼저",
            "  구현",
        ]),
        ("오늘 산출물", GREEN, [
            "노트 작성 & 홈 화면",
            "",
            "• note/create.tsx",
            "• (tabs)/index.tsx",
            "• services/notes.ts",
            "• types/index.ts",
            "• 하단 탭바 구성",
            "  홈 / 캘린더 / 설정",
        ]),
        ("저장 위치", ORANGE, [
            "📁 app/note/",
            "└── create.tsx",
            "",
            "📁 app/(tabs)/",
            "└── index.tsx",
            "",
            "📁 services/",
            "└── notes.ts",
        ]),
    ]
)

# ════════════════════════════════════════
# 슬라이드 7 — 다음 계획
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DARK)

add_textbox(slide, "다음 개발 계획", 0.5, 0.3, 12, 0.7,
            font_size=30, bold=True, color=WHITE)

plans = [
    (GREEN,  "3단계 마무리",  ["노트 상세 화면", "자가 체크 기능", "복습 날짜 자동 계산 로직"]),
    (ORANGE, "4단계",         ["캘린더 화면 (월간 뷰)", "스트릭 기능", "연속 복습 일수 표시"]),
    (RGBColor(0xE7,0x4C,0x3C), "5단계", ["Claude API 연동", "노트 기반 문제 자동 생성", "AI 테스트 화면"]),
    (BLUE,   "6단계",         ["설정 화면", "복습 알림 기능", "전체 UI 다듬기"]),
]

for i, (color, title, items) in enumerate(plans):
    x = 0.4 + i * 3.1
    box = slide.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(2.85), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE
    for item in items:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "\n• " + item
        r2.font.name = FONT_BODY
        r2.font.size = Pt(13)
        r2.font.color.rgb = WHITE

add_textbox(slide, "프로토타입 목표: 1~4단계 완성 후 실제 폰 시연",
            0.5, 6.4, 12, 0.5,
            font_size=14, color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.CENTER)

# ── 저장 ──────────────────────────────────
output = r"C:\Users\jisoo\OneDrive\Desktop\클로드 바이브 코딩\바이브 코딩 개인 노트\망각곡선\망각곡선_앱_발표.pptx"
prs.save(output)
print("저장 완료: " + output)
